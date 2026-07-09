"""
PPTX -> draft slide JSON. Mirrors the manual extraction process already
proven for Modul 2 (see project memory): read text via text_frame, read
tables via shape.table (not flat text, to keep nested header/subtotal rows
correct), and pull out embedded pictures (including "hidden" screenshots of
formulas/tables that have no text_frame) as base64 data URIs.
"""
import base64
from pptx import Presentation
from pptx.util import Emu


def _shape_text(shape):
    if not shape.has_text_frame:
        return ''
    lines = []
    for para in shape.text_frame.paragraphs:
        text = ''.join(run.text for run in para.runs)
        if text.strip():
            lines.append(text)
    return '\n'.join(lines)


def _shape_table(shape):
    if not shape.has_table:
        return None
    table = shape.table
    rows = []
    for row in table.rows:
        rows.append([cell.text for cell in row.cells])
    return rows


def _shape_image(shape):
    try:
        image = shape.image
    except Exception:
        return None
    ext = image.ext or 'png'
    mime = 'jpeg' if ext.lower() in ('jpg', 'jpeg') else ext.lower()
    b64 = base64.b64encode(image.blob).decode('ascii')
    return f'data:image/{mime};base64,{b64}'


def extract(path):
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        tables = []
        images = []
        for shape in slide.shapes:
            if getattr(shape, 'has_table', False) and shape.has_table:
                t = _shape_table(shape)
                if t:
                    tables.append(t)
                continue
            if getattr(shape, 'has_text_frame', False) and shape.has_text_frame:
                t = _shape_text(shape)
                if t:
                    texts.append(t)
            if shape.shape_type == 13 or (hasattr(shape, 'image') and not getattr(shape, 'has_text_frame', False)):
                img = _shape_image(shape)
                if img:
                    images.append(img)
        slides.append({
            'slideNo': i,
            'texts': texts,
            'tables': tables,
            'images': images,
        })
    return slides
